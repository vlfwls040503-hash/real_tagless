from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# === 색상 팔레트 ===
BG_DARK = RGBColor(0x1B, 0x1B, 0x2F)
BG_SECTION = RGBColor(0x22, 0x22, 0x3A)
ACCENT_BLUE = RGBColor(0x4F, 0xC3, 0xF7)
ACCENT_ORANGE = RGBColor(0xFF, 0xB7, 0x4D)
ACCENT_GREEN = RGBColor(0x66, 0xBB, 0x6A)
ACCENT_RED = RGBColor(0xEF, 0x53, 0x50)
ACCENT_PURPLE = RGBColor(0xCE, 0x93, 0xD8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBD, 0xBD, 0xBD)
MEDIUM_GRAY = RGBColor(0x90, 0x90, 0xA0)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="맑은 고딕"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_text(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "맑은 고딕"
        p.space_after = spacing
    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ============================================================
# 슬라이드 1: 표지
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_accent_line(slide, Inches(1.5), Inches(2.5), Inches(10.3), ACCENT_BLUE)
add_text(slide, Inches(1.5), Inches(2.7), Inches(10.3), Inches(1.5),
         "태그리스 전용 게이트 분리 배치의\n통행비용 절감 효과 분석",
         font_size=38, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
add_text(slide, Inches(1.5), Inches(4.3), Inches(10.3), Inches(0.6),
         "성수역 보행 시뮬레이션 기반",
         font_size=22, color=ACCENT_BLUE, bold=False, alignment=PP_ALIGN.LEFT)
add_text(slide, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.8),
         "교통공학과  |  졸업설계  |  2026",
         font_size=18, color=MEDIUM_GRAY, alignment=PP_ALIGN.LEFT)


# ============================================================
# 슬라이드 2: 태그리스 시스템이란?
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
         "태그리스 시스템이란?", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(3), ACCENT_BLUE)

add_bullet_text(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(3.5), [
    "블루투스(BLE) 기반, 카드 태그 없이 게이트 통과 시 자동 결제",
    "고속도로 하이패스와 유사한 원리",
    "2023년 우이신설선에서 세계 최초 상용화",
    "서울교통공사 1~8호선 시범사업 확대 중",
], font_size=18, color=LIGHT_GRAY)

# 비교 박스
box1 = add_shape(slide, Inches(7), Inches(1.5), Inches(2.7), Inches(2.2), RGBColor(0x2A, 0x2A, 0x45), ACCENT_ORANGE)
add_text(slide, Inches(7.2), Inches(1.65), Inches(2.3), Inches(0.5),
         "기존 태그 방식", font_size=16, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(7.2), Inches(2.15), Inches(2.3), Inches(1.3),
         "카드/스마트폰을\n단말기에 접촉\n(NFC, ~20cm)",
         font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

box2 = add_shape(slide, Inches(10.1), Inches(1.5), Inches(2.7), Inches(2.2), RGBColor(0x2A, 0x2A, 0x45), ACCENT_BLUE)
add_text(slide, Inches(10.3), Inches(1.65), Inches(2.3), Inches(0.5),
         "태그리스 방식", font_size=16, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(10.3), Inches(2.15), Inches(2.3), Inches(1.3),
         "접촉 없이\n게이트 통과만으로 결제\n(BLE, ~10m)",
         font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# 하단 강조
add_shape(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.5), RGBColor(0x1A, 0x33, 0x50), ACCENT_BLUE)
add_text(slide, Inches(1.2), Inches(5.4), Inches(11), Inches(1.1),
         "핵심: 태그리스는 '멈추지 않고 통과'할 수 있어 게이트 처리 속도가 근본적으로 다름",
         font_size=18, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.LEFT)


# ============================================================
# 슬라이드 3: 현행 운영 방식과 서비스 시간 차이 (합침)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
         "현행 게이트 운영과 서비스 시간 차이", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(3), ACCENT_BLUE)

# 좌측: 현행 운영
add_shape(slide, Inches(0.8), Inches(1.4), Inches(5.8), Inches(2.6), RGBColor(0x22, 0x22, 0x3A), ACCENT_BLUE)
add_text(slide, Inches(1.1), Inches(1.55), Inches(5.2), Inches(0.5),
         "현행 운영 방식", font_size=20, color=ACCENT_BLUE, bold=True)
add_bullet_text(slide, Inches(1.1), Inches(2.1), Inches(5.2), Inches(1.7), [
    "대부분 태그(NFC) 전용 게이트",
    "역당 1~2대만 태그/태그리스 겸용 운영",
    "겸용 게이트에 두 유형의 이용자가 같은 줄에 혼재",
], font_size=15, color=LIGHT_GRAY, spacing=Pt(6))

# 우측: 서비스 시간
add_shape(slide, Inches(7), Inches(1.4), Inches(5.5), Inches(2.6), RGBColor(0x22, 0x22, 0x3A), ACCENT_ORANGE)
add_text(slide, Inches(7.3), Inches(1.55), Inches(5), Inches(0.5),
         "서비스 시간 (= 1명 처리 시간)", font_size=20, color=ACCENT_ORANGE, bold=True)
add_bullet_text(slide, Inches(7.3), Inches(2.1), Inches(5), Inches(1.7), [
    "태그 이용자: 정지 > 탭 > 문 열림 > 통과  (2~3초)",
    "태그리스 이용자: 감속 없이 연속 통과  (1초 이하)",
    "같은 게이트에서 서비스 시간이 2~3배 차이",
], font_size=15, color=LIGHT_GRAY, spacing=Pt(6))

# 게이트 배치 다이어그램
add_text(slide, Inches(0.8), Inches(4.3), Inches(3), Inches(0.4),
         "게이트 배치 예시", font_size=14, color=MEDIUM_GRAY, bold=True)

y_base = Inches(4.7)
gate_w = Inches(0.7)
gate_h = Inches(0.8)
gap = Inches(0.12)

for i in range(10):
    x = Inches(0.8) + i * (gate_w + gap)
    if i >= 8:
        color = RGBColor(0x1A, 0x44, 0x66)
        border = ACCENT_BLUE
        label_color = ACCENT_BLUE
        label = "겸용"
    else:
        color = RGBColor(0x3A, 0x2A, 0x1A)
        border = ACCENT_ORANGE
        label_color = ACCENT_ORANGE
        label = "태그"
    add_shape(slide, x, y_base, gate_w, gate_h, color, border)
    add_text(slide, x, y_base + Inches(0.2), gate_w, Inches(0.4),
             label, font_size=11, color=label_color, bold=True, alignment=PP_ALIGN.CENTER)

# 하단 강조
add_shape(slide, Inches(0.8), Inches(5.9), Inches(11.7), Inches(1.0), RGBColor(0x40, 0x1A, 0x1A), ACCENT_RED)
add_text(slide, Inches(1.2), Inches(6.05), Inches(11), Inches(0.7),
         "빠른 이용자(태그리스)가 느린 이용자(태그) 뒤에서 불필요하게 대기 => 태그리스의 연속 통과 이점 상실",
         font_size=17, color=ACCENT_RED, bold=True)


# ============================================================
# 슬라이드 4: 선행연구 (합침)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
         "선행연구", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(3), ACCENT_GREEN)

# 연구 1: 대기행렬 이론
add_shape(slide, Inches(0.8), Inches(1.4), Inches(3.7), Inches(2.5), RGBColor(0x1A, 0x2A, 0x1A), ACCENT_GREEN)
add_text(slide, Inches(1.0), Inches(1.55), Inches(3.3), Inches(0.4),
         "이질적 서비스 시간 대기행렬", font_size=14, color=ACCENT_GREEN, bold=True)
add_bullet_text(slide, Inches(1.0), Inches(2.0), Inches(3.3), Inches(1.7), [
    "서비스 시간이 다른 고객이 혼재하면",
    "전체 대기시간 증가 (IEEE, 2023)",
    "",
    "게이트 서비스 시간 = 역사 용량의",
    "핵심 변수 (TRC, 2013)",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(3))

# 연구 2: 혼합류
add_shape(slide, Inches(4.8), Inches(1.4), Inches(3.7), Inches(2.5), RGBColor(0x1A, 0x2A, 0x1A), ACCENT_GREEN)
add_text(slide, Inches(5.0), Inches(1.55), Inches(3.3), Inches(0.4),
         "혼합류 통과 효율 저하", font_size=14, color=ACCENT_GREEN, bold=True)
add_bullet_text(slide, Inches(5.0), Inches(2.0), Inches(3.3), Inches(1.7), [
    "통과 특성이 다른 보행자가 혼합되면",
    "게이트 통과 효율 저하",
    "(Safety Science, 2024)",
    "",
    "혼합 비율이 높을수록 효과 뚜렷",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(3))

# 연구 3: 분리 효과 + 배치
add_shape(slide, Inches(8.8), Inches(1.4), Inches(3.7), Inches(2.5), RGBColor(0x1A, 0x2A, 0x1A), ACCENT_GREEN)
add_text(slide, Inches(9.0), Inches(1.55), Inches(3.3), Inches(0.4),
         "전용 분리 및 배치 효과", font_size=14, color=ACCENT_GREEN, bold=True)
add_bullet_text(slide, Inches(9.0), Inches(2.0), Inches(3.3), Inches(1.7), [
    "처리 속도가 다른 이용자를 전용",
    "레인으로 분리 시 처리량 증가",
    "(J. of Trans. Security, 2022)",
    "",
    "게이트 배치 변경만으로도",
    "보행류 효율 변화 (Tanaka, 2022)",
], font_size=13, color=LIGHT_GRAY, spacing=Pt(3))

# 시사점
add_shape(slide, Inches(0.8), Inches(4.4), Inches(11.7), Inches(2.5), RGBColor(0x22, 0x22, 0x3A))
add_text(slide, Inches(1.1), Inches(4.6), Inches(11), Inches(0.4),
         "시사점", font_size=18, color=ACCENT_GREEN, bold=True)
add_bullet_text(slide, Inches(1.1), Inches(5.1), Inches(11), Inches(1.6), [
    "서비스 시간이 다른 이용자가 섞이면 대기시간 증가 => 이론적으로 검증됨",
    "통과 특성이 다른 보행자 혼합 시 게이트 효율 저하 => 실험적으로 검증됨",
    "전용 레인 분리 및 배치 변경만으로 전체 처리량 개선 가능 => 타 분야에서 검증됨",
    "=> 태그(2~3초)와 태그리스(1초 이하)가 섞인 겸용 게이트에도 동일한 문제·해법 적용 가능",
], font_size=15, color=LIGHT_GRAY, spacing=Pt(6))


# ============================================================
# 슬라이드 5: 연구 공백
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
         "연구 공백", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(3), ACCENT_RED)

# 있는 것: 태그리스 게이트의 효율성
add_shape(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(4.0), RGBColor(0x1A, 0x2A, 0x1A), ACCENT_GREEN)
add_text(slide, Inches(1.1), Inches(1.55), Inches(5), Inches(0.5),
         "확인된 것", font_size=18, color=ACCENT_GREEN, bold=True)
add_bullet_text(slide, Inches(1.1), Inches(2.1), Inches(5), Inches(3.1), [
    "태그리스 게이트는 서비스 시간이 1초 이하로",
    "기존 태그(2~3초) 대비 처리 효율이 월등히 높음",
    "",
    "서비스 시간이 다른 이용자 혼재 시 대기시간 증가",
    "전용 레인 분리 시 전체 처리량 개선",
    "게이트 배치 변경만으로 보행류 효율 변화 가능",
    "",
    "=> 태그리스 전용 게이트 분리 배치는",
    "   이론적으로 통행비용 절감이 기대됨",
], font_size=14, color=LIGHT_GRAY, spacing=Pt(4))

# 없는 것
add_shape(slide, Inches(7), Inches(1.4), Inches(5.5), Inches(4.0), RGBColor(0x40, 0x1A, 0x1A), ACCENT_RED)
add_text(slide, Inches(7.3), Inches(1.55), Inches(5), Inches(0.5),
         "아직 연구되지 않은 것", font_size=18, color=ACCENT_RED, bold=True)
add_bullet_text(slide, Inches(7.3), Inches(2.1), Inches(5), Inches(3.1), [
    "지하철 게이트에서 태그/태그리스 혼재 시",
    "병목에 대한 정량적 분석",
    "",
    "겸용 게이트를 전용으로 분리 전환했을 때",
    "실제 통행비용 절감 효과",
    "",
    "태그리스 비율별 전용 전환의 임계점",
    "(몇 % 이상이면 분리가 유의미한지)",
], font_size=14, color=LIGHT_GRAY, spacing=Pt(4))

# 하단 강조
add_shape(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.2), RGBColor(0x22, 0x22, 0x3A))
add_text(slide, Inches(1.2), Inches(5.95), Inches(11), Inches(0.9),
         "태그리스의 효율성은 확인되었으나, 전용 분리 배치의 정량적 효과는 검증된 바 없음\n=> 본 연구의 필요성",
         font_size=20, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# 슬라이드 6: 연구 목적
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
         "연구 목적", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(3), ACCENT_BLUE)

purposes = [
    ("01", "통행비용 비교", "현행 겸용 게이트 운영과 태그리스 전용 게이트\n분리 배치의 통행비용을 정량적으로 비교", ACCENT_BLUE),
    ("02", "임계점 도출", "태그리스 이용자 비율이 몇 % 이상일 때\n전용 게이트 전환이 유의미한 효과를 갖는지 도출", ACCENT_ORANGE),
    ("03", "정책 근거 제시", "추가 인프라 변경 없이 게이트 설정 변경만으로\n달성 가능한 통행비용 절감 효과 제시", ACCENT_GREEN),
]

for i, (num, title, desc, color) in enumerate(purposes):
    y = Inches(1.6) + i * Inches(1.8)
    add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(1.5), RGBColor(0x22, 0x22, 0x3A), color)
    add_text(slide, Inches(1.2), y + Inches(0.15), Inches(1), Inches(0.5),
             num, font_size=28, color=color, bold=True)
    add_text(slide, Inches(2.3), y + Inches(0.15), Inches(3), Inches(0.5),
             title, font_size=20, color=WHITE, bold=True)
    add_text(slide, Inches(2.3), y + Inches(0.7), Inches(9.5), Inches(0.7),
             desc, font_size=15, color=LIGHT_GRAY)


# ============================================================
# 슬라이드 7: 연구 방법 개요
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
         "연구 방법 개요", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(3), ACCENT_BLUE)

# 대상역
add_shape(slide, Inches(0.8), Inches(1.5), Inches(3.5), Inches(3.0), RGBColor(0x22, 0x22, 0x3A), ACCENT_BLUE)
add_text(slide, Inches(1.0), Inches(1.65), Inches(3.1), Inches(0.5),
         "대상역: 성수역", font_size=18, color=ACCENT_BLUE, bold=True)
add_bullet_text(slide, Inches(1.0), Inches(2.2), Inches(3.1), Inches(2.0), [
    "2호선, 환승 없음",
    "구조 단순 => 시뮬레이션 용이",
    "승객/게이트 비율 상위 3위",
    "  (2,707명/게이트)",
    "일평균 105,574명",
], font_size=14, color=LIGHT_GRAY, spacing=Pt(4))

# 시뮬레이션
add_shape(slide, Inches(4.6), Inches(1.5), Inches(3.5), Inches(3.0), RGBColor(0x22, 0x22, 0x3A), ACCENT_ORANGE)
add_text(slide, Inches(4.8), Inches(1.65), Inches(3.1), Inches(0.5),
         "시뮬레이션", font_size=18, color=ACCENT_ORANGE, bold=True)
add_bullet_text(slide, Inches(4.8), Inches(2.2), Inches(3.1), Inches(2.0), [
    "Social Force Model 기반",
    "보행 시뮬레이션",
    "게이트 대기열 + 서비스시간",
    "차이를 반영한 모델링",
], font_size=14, color=LIGHT_GRAY, spacing=Pt(4))

# 시나리오
add_shape(slide, Inches(8.4), Inches(1.5), Inches(4.1), Inches(3.0), RGBColor(0x22, 0x22, 0x3A), ACCENT_GREEN)
add_text(slide, Inches(8.6), Inches(1.65), Inches(3.7), Inches(0.5),
         "시나리오 설계", font_size=18, color=ACCENT_GREEN, bold=True)
add_bullet_text(slide, Inches(8.6), Inches(2.2), Inches(3.7), Inches(2.0), [
    "기본: 현행 겸용 게이트 운영",
    "개선: 태그리스 전용 게이트 분리",
    "태그리스 비율: 10~80% (5단계)",
    "비율별 기본/개선 각각 시뮬레이션",
], font_size=14, color=LIGHT_GRAY, spacing=Pt(4))

# 평가지표
add_shape(slide, Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.8), RGBColor(0x22, 0x22, 0x3A))
add_text(slide, Inches(1.1), Inches(5.15), Inches(3), Inches(0.5),
         "평가 지표", font_size=18, color=ACCENT_PURPLE, bold=True)

metrics = ["평균 대기시간", "최대 대기행렬 길이", "총 통과시간\n(게이트 진입~계단 도달)", "보행자 밀도 / LOS"]
for i, m in enumerate(metrics):
    x = Inches(1.1) + i * Inches(2.9)
    add_shape(slide, x, Inches(5.7), Inches(2.6), Inches(0.9), RGBColor(0x2A, 0x2A, 0x45), ACCENT_PURPLE)
    add_text(slide, x + Inches(0.1), Inches(5.8), Inches(2.4), Inches(0.7),
             m, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# 슬라이드 8: 기대 결과 및 활용
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
         "기대 결과 및 활용", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(3), ACCENT_BLUE)

results = [
    ("기본 vs 개선 시나리오 비교", "태그리스 비율별 대기시간, 통과시간 차이를 정량화\n비율이 높아질수록 전용 게이트 전환 효과 증가 예상", ACCENT_BLUE),
    ("전용 전환 임계점 도출", "\"태그리스 비율 X% 이상에서 전용 게이트 전환이 유의미\"\n정책 결정을 위한 정량적 기준 제시", ACCENT_ORANGE),
    ("정책 제언", "현 인프라에서 게이트 설정 변경만으로 달성 가능한 효과 제시\n서울교통공사 태그리스 확대 시 게이트 운영 전략 근거", ACCENT_GREEN),
]

for i, (title, desc, color) in enumerate(results):
    y = Inches(1.6) + i * Inches(1.8)
    add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(1.5), RGBColor(0x22, 0x22, 0x3A), color)
    add_text(slide, Inches(1.2), y + Inches(0.15), Inches(10), Inches(0.5),
             title, font_size=20, color=color, bold=True)
    add_text(slide, Inches(1.2), y + Inches(0.7), Inches(10.5), Inches(0.7),
             desc, font_size=15, color=LIGHT_GRAY)


# ============================================================
# 슬라이드 9: 향후 일정 + 감사합니다
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
         "향후 일정", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.05), Inches(3), ACCENT_BLUE)

schedule = [
    ("1단계", "성수역 평면도 확보 및 공간 모델링", "3~4월"),
    ("2단계", "시뮬레이션 모델 구축 및 파라미터 설정", "4~5월"),
    ("3단계", "시나리오별 시뮬레이션 실행", "5~6월"),
    ("4단계", "결과 분석 및 논문 작성", "6~7월"),
]

colors = [ACCENT_BLUE, ACCENT_ORANGE, ACCENT_GREEN, ACCENT_PURPLE]

for i, (phase, desc, period) in enumerate(schedule):
    y = Inches(1.5) + i * Inches(1.35)
    add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(1.1), RGBColor(0x22, 0x22, 0x3A), colors[i])
    add_text(slide, Inches(1.2), y + Inches(0.1), Inches(1.5), Inches(0.5),
             phase, font_size=18, color=colors[i], bold=True)
    add_text(slide, Inches(2.8), y + Inches(0.1), Inches(6.5), Inches(0.5),
             desc, font_size=17, color=WHITE)
    add_text(slide, Inches(10), y + Inches(0.1), Inches(2.2), Inches(0.5),
             period, font_size=17, color=MEDIUM_GRAY, alignment=PP_ALIGN.RIGHT)

# 하단
add_text(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5),
         "감사합니다", font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# 저장
# ============================================================
output_path = r"C:\Users\aaron\Desktop\tagless_presentation_v2.pptx"
prs.save(output_path)
print(f"저장 완료: {output_path}")
